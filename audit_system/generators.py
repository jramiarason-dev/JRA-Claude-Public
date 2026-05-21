"""Output generators: PowerPoint, Excel, Word documents for audit agents."""

import os
from datetime import datetime
from pathlib import Path

# ── PowerPoint ────────────────────────────────────────────────────────────────

def generate_audit_plan_ppt(data: dict, output_dir: str = "outputs") -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY   = RGBColor(0x0D, 0x2B, 0x55)
    BLUE   = RGBColor(0x1A, 0x5F, 0x9E)
    GOLD   = RGBColor(0xC8, 0xA0, 0x32)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)

    RISK_COLORS = {
        "Critical": RGBColor(0xC0, 0x39, 0x2B),
        "High":     RGBColor(0xE7, 0x4C, 0x3C),
        "Medium":   RGBColor(0xE6, 0x7E, 0x22),
        "Low":      RGBColor(0x27, 0xAE, 0x60),
    }

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # completely blank layout

    def add_rect(slide, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background() if line_color is None else None
        return shape

    def add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                     color=WHITE, align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txb

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 0.3, 7.5, GOLD)
    add_rect(slide, 0, 5.0, 13.33, 0.05, GOLD)

    audit_title = data.get("audit_title", "Internal Audit Plan")
    add_text_box(slide, "INTERNAL AUDIT", 0.5, 1.5, 12, 0.8, 14, bold=False, color=GOLD)
    add_text_box(slide, audit_title, 0.5, 2.2, 12, 1.5, 32, bold=True, color=WHITE)
    add_text_box(slide, data.get("audit_scope", ""), 0.5, 3.8, 12, 1.0, 14, color=RGBColor(0xBB, 0xCC, 0xDD))

    now = datetime.now().strftime("%B %Y")
    add_text_box(slide, f"Prepared: {now}  |  CONFIDENTIAL", 0.5, 6.8, 12, 0.5, 10, color=GOLD, align=PP_ALIGN.CENTER)

    # ── Slide 2: Scope & Objectives ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_rect(slide, 0, 1.1, 13.33, 0.05, GOLD)
    add_text_box(slide, "AUDIT SCOPE & OBJECTIVES", 0.4, 0.2, 12, 0.8, 20, bold=True, color=WHITE)

    objectives = data.get("objectives", [])
    y = 1.4
    if data.get("audit_scope"):
        add_text_box(slide, "Scope", 0.4, y, 12, 0.4, 13, bold=True, color=NAVY)
        add_text_box(slide, data["audit_scope"], 0.4, y + 0.35, 12.5, 0.8, 11, color=RGBColor(0x33, 0x33, 0x33))
        y += 1.3

    if objectives:
        add_text_box(slide, "Key Objectives", 0.4, y, 12, 0.4, 13, bold=True, color=NAVY)
        y += 0.4
        for obj in objectives[:6]:
            add_text_box(slide, f"▸  {obj}", 0.6, y, 12, 0.35, 11, color=RGBColor(0x22, 0x22, 0x22))
            y += 0.35

    # Jurisdiction pills
    jurisdictions = ["Switzerland (FINMA)", "Singapore (MAS)", "Hong Kong (SFC/HKMA)",
                     "Bahamas (SCB)", "European Union", "United Kingdom (FCA/PRA)"]
    add_text_box(slide, "Jurisdictions in Scope", 0.4, 6.0, 12, 0.4, 12, bold=True, color=NAVY)
    x_j = 0.4
    for j in jurisdictions:
        w = len(j) * 0.095 + 0.2
        add_rect(slide, x_j, 6.45, w, 0.32, BLUE)
        add_text_box(slide, j, x_j + 0.05, 6.47, w, 0.28, 9, color=WHITE)
        x_j += w + 0.12

    # ── Slides 3+: One per audit subject ──────────────────────────────────────
    subjects = data.get("subjects", [])
    for i, subj in enumerate(subjects, start=1):
        slide = prs.slides.add_slide(blank)
        risk  = subj.get("risk_level", "Medium")
        rcol  = RISK_COLORS.get(risk, RISK_COLORS["Medium"])

        add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
        add_rect(slide, 0, 1.1, 13.33, 0.05, rcol)
        add_text_box(slide, f"AUDIT SUBJECT {i:02d}", 0.4, 0.15, 10, 0.4, 11, color=GOLD)
        add_text_box(slide, subj.get("subject_name", ""), 0.4, 0.5, 11, 0.55, 22, bold=True, color=WHITE)

        # Risk badge
        add_rect(slide, 11.3, 0.2, 1.8, 0.55, rcol)
        add_text_box(slide, f"⚠ {risk.upper()} RISK", 11.3, 0.27, 1.8, 0.4, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, subj.get("description", ""), 0.4, 1.3, 12.5, 0.9, 11, color=RGBColor(0x33, 0x33, 0x33))

        # Objectives column
        obj_list = subj.get("key_objectives", [])
        add_rect(slide, 0.3, 2.35, 6.2, 0.35, BLUE)
        add_text_box(slide, "KEY OBJECTIVES", 0.4, 2.38, 6.0, 0.3, 10, bold=True, color=WHITE)
        y = 2.75
        for obj in obj_list[:5]:
            add_text_box(slide, f"• {obj}", 0.4, y, 6.0, 0.35, 10, color=RGBColor(0x22, 0x22, 0x22))
            y += 0.35

        # Regulatory refs column
        reg_refs = subj.get("regulatory_references", [])
        add_rect(slide, 6.8, 2.35, 6.2, 0.35, BLUE)
        add_text_box(slide, "REGULATORY REFERENCES", 6.9, 2.38, 6.0, 0.3, 10, bold=True, color=WHITE)
        y = 2.75
        for ref in reg_refs[:6]:
            add_text_box(slide, f"⬡ {ref}", 6.9, y, 6.0, 0.32, 10, color=RGBColor(0x22, 0x22, 0x22))
            y += 0.32

        # Jurisdictions
        jurs = subj.get("applicable_jurisdictions", [])
        if jurs:
            add_text_box(slide, "APPLICABLE JURISDICTIONS: " + "  ·  ".join(jurs), 0.4, 6.9, 12.5, 0.35, 10, color=NAVY, bold=True)

    # ── Summary slide: Risk overview ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_rect(slide, 0, 1.1, 13.33, 0.05, GOLD)
    add_text_box(slide, "AUDIT RISK OVERVIEW", 0.4, 0.25, 12, 0.7, 22, bold=True, color=WHITE)

    from collections import Counter
    risk_counts = Counter(s.get("risk_level", "Medium") for s in subjects)
    boxes = [("Critical", RISK_COLORS["Critical"]), ("High", RISK_COLORS["High"]),
             ("Medium", RISK_COLORS["Medium"]), ("Low", RISK_COLORS["Low"])]
    x_b = 1.0
    for label, col in boxes:
        add_rect(slide, x_b, 1.5, 2.5, 1.6, col)
        add_text_box(slide, str(risk_counts.get(label, 0)), x_b, 1.7, 2.5, 1.0, 48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x_b, 2.75, 2.5, 0.4, 13, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
        x_b += 2.8

    y = 3.4
    for subj in subjects:
        risk  = subj.get("risk_level", "Medium")
        rcol  = RISK_COLORS.get(risk, RISK_COLORS["Medium"])
        add_rect(slide, 0.3, y, 0.25, 0.3, rcol)
        add_text_box(slide, subj.get("subject_name", ""), 0.7, y, 9, 0.3, 11, color=RGBColor(0x22, 0x22, 0x22))
        add_text_box(slide, risk, 10.0, y, 2.5, 0.3, 11, bold=True, color=rcol)
        y += 0.33

    # ── Save ──────────────────────────────────────────────────────────────────
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename  = f"{output_dir}/Audit_Plan_{timestamp}.pptx"
    prs.save(filename)
    return filename


# ── Excel ─────────────────────────────────────────────────────────────────────

def generate_audit_procedures_excel(data: dict, output_dir: str = "outputs") -> str:
    import openpyxl
    from openpyxl.styles import (PatternFill, Font, Alignment, Border, Side,
                                  GradientFill)
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()

    NAVY_HEX  = "0D2B55"
    BLUE_HEX  = "1A5F9E"
    GOLD_HEX  = "C8A032"
    LGRAY_HEX = "F0F4F8"
    WHITE_HEX = "FFFFFF"

    RISK_FILLS = {
        "Critical": PatternFill("solid", fgColor="C0392B"),
        "High":     PatternFill("solid", fgColor="E74C3C"),
        "Medium":   PatternFill("solid", fgColor="E67E22"),
        "Low":      PatternFill("solid", fgColor="27AE60"),
    }
    RISK_FONTS = {k: Font(bold=True, color=WHITE_HEX) for k in RISK_FILLS}

    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    def header_style(cell, bg=NAVY_HEX, font_size=11):
        cell.fill   = PatternFill("solid", fgColor=bg)
        cell.font   = Font(bold=True, color=WHITE_HEX, size=font_size)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border

    def data_style(cell, wrap=True, bold=False, bg=None):
        cell.font   = Font(bold=bold, size=10)
        cell.alignment = Alignment(vertical="top", wrap_text=wrap)
        cell.border = border
        if bg:
            cell.fill = PatternFill("solid", fgColor=bg)

    # ── Cover sheet ───────────────────────────────────────────────────────────
    ws_cover = wb.active
    ws_cover.title = "Cover"
    ws_cover.column_dimensions["A"].width = 3
    ws_cover.column_dimensions["B"].width = 50
    ws_cover.column_dimensions["C"].width = 30

    ws_cover["B2"] = "INTERNAL AUDIT PROCEDURES"
    ws_cover["B2"].font = Font(bold=True, size=22, color=NAVY_HEX)
    ws_cover["B3"] = data.get("audit_title", "Audit Programme")
    ws_cover["B3"].font = Font(bold=True, size=16, color=BLUE_HEX)
    ws_cover["B5"] = f"Date: {datetime.now().strftime('%d %B %Y')}"
    ws_cover["B5"].font = Font(size=12, color="555555")
    ws_cover["B6"] = "Classification: CONFIDENTIAL"
    ws_cover["B6"].font = Font(size=12, color="AA0000", bold=True)

    ws_cover["B8"] = "Jurisdictions in Scope"
    ws_cover["B8"].font = Font(bold=True, size=12, color=NAVY_HEX)
    jurs = ["Switzerland (FINMA)", "Singapore (MAS)", "Hong Kong (SFC/HKMA)",
            "Bahamas (SCB)", "European Union", "United Kingdom (FCA/PRA)"]
    for r, j in enumerate(jurs, start=9):
        ws_cover[f"B{r}"] = f"  ▸  {j}"
        ws_cover[f"B{r}"].font = Font(size=11)

    # Subject index
    ws_cover["B16"] = "Audit Subjects"
    ws_cover["B16"].font = Font(bold=True, size=12, color=NAVY_HEX)
    subjects = data.get("subjects", [])
    for r, subj in enumerate(subjects, start=17):
        ws_cover[f"B{r}"] = f"  {r-16:02d}.  {subj.get('subject_name', '')}"
        ws_cover[f"C{r}"] = subj.get("procedures", [{}])[0].get("risk_level", "") if subj.get("procedures") else ""
        risk = ws_cover[f"C{r}"].value
        if risk in RISK_FILLS:
            ws_cover[f"C{r}"].fill = RISK_FILLS[risk]
            ws_cover[f"C{r}"].font = RISK_FONTS[risk]
            ws_cover[f"C{r}"].alignment = Alignment(horizontal="center")

    # ── One sheet per subject ─────────────────────────────────────────────────
    COLS = [
        ("Ref",         8),
        ("Objective",   30),
        ("Procedure",   40),
        ("Test Steps",  45),
        ("Sample",      12),
        ("Evidence",    35),
        ("Risks",       35),
        ("Risk Level",  12),
        ("Control",     30),
        ("Reg. Refs",   30),
    ]

    for subj in subjects:
        sheet_name = subj.get("subject_name", "Subject")[:31]
        ws = wb.create_sheet(title=sheet_name)

        # Title row
        ws.merge_cells("A1:J1")
        title_cell = ws["A1"]
        title_cell.value   = f"{data.get('audit_title', '')}  ·  {sheet_name}"
        title_cell.fill    = PatternFill("solid", fgColor=NAVY_HEX)
        title_cell.font    = Font(bold=True, size=14, color=WHITE_HEX)
        title_cell.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 30

        # Header row
        for col_i, (col_name, col_w) in enumerate(COLS, start=1):
            cell = ws.cell(row=2, column=col_i, value=col_name)
            header_style(cell, bg=BLUE_HEX)
            ws.column_dimensions[get_column_letter(col_i)].width = col_w
        ws.row_dimensions[2].height = 25

        # Data rows
        procedures = subj.get("procedures", [])
        for row_i, proc in enumerate(procedures, start=3):
            risk = proc.get("risk_level", "Medium")
            bg   = LGRAY_HEX if row_i % 2 == 0 else WHITE_HEX

            values = [
                proc.get("ref", f"P{row_i-2:03d}"),
                proc.get("objective", ""),
                proc.get("procedure", ""),
                "\n".join(f"  {j+1}. {s}" for j, s in enumerate(proc.get("test_steps", []))),
                proc.get("sample_size", ""),
                proc.get("evidence_required", ""),
                "\n".join(f"• {r}" for r in proc.get("associated_risks", [])),
                risk,
                proc.get("control_objective", ""),
                "\n".join(proc.get("regulatory_refs", [])),
            ]

            for col_i, val in enumerate(values, start=1):
                cell = ws.cell(row=row_i, column=col_i, value=val)
                if col_i == 8:  # Risk Level column
                    cell.fill  = RISK_FILLS.get(risk, PatternFill())
                    cell.font  = RISK_FONTS.get(risk, Font())
                    cell.alignment = Alignment(horizontal="center", vertical="center")
                    cell.border = border
                else:
                    data_style(cell, bg=bg)

            ws.row_dimensions[row_i].height = max(
                60, len(proc.get("test_steps", [])) * 15
            )

        ws.freeze_panes = "A3"

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename  = f"{output_dir}/Audit_Procedures_{timestamp}.xlsx"
    wb.save(filename)
    return filename


# ── Word Document ─────────────────────────────────────────────────────────────

def generate_audit_report_docx(data: dict, output_dir: str = "outputs") -> str:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor as DocxRGB, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    NAVY  = DocxRGB(0x0D, 0x2B, 0x55)
    BLUE  = DocxRGB(0x1A, 0x5F, 0x9E)
    GOLD  = DocxRGB(0xC8, 0xA0, 0x32)
    WHITE = DocxRGB(0xFF, 0xFF, 0xFF)

    RISK_COLORS_DOCX = {
        "Critical":     DocxRGB(0xC0, 0x39, 0x2B),
        "High":         DocxRGB(0xE7, 0x4C, 0x3C),
        "Medium":       DocxRGB(0xE6, 0x7E, 0x22),
        "Low":          DocxRGB(0x27, 0xAE, 0x60),
        "Informational":DocxRGB(0x29, 0x80, 0xB9),
    }

    RISK_HEX = {
        "Critical":     "C0392B",
        "High":         "E74C3C",
        "Medium":       "E67E22",
        "Low":          "27AE60",
        "Informational":"2980B9",
    }

    def set_cell_bg(cell, hex_color):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), hex_color)
        tcPr.append(shd)

    def add_heading(text, level=1, color=None):
        p = doc.add_paragraph()
        p.style = doc.styles["Heading 1"] if level == 1 else doc.styles["Heading 2"]
        run = p.add_run(text)
        run.font.size = Pt(16 if level == 1 else 13)
        run.font.bold = True
        run.font.color.rgb = color or NAVY
        return p

    def add_body(text, bold=False, italic=False, color=None, size=10.5):
        p   = doc.add_paragraph()
        run = p.add_run(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        p.paragraph_format.space_after = Pt(4)
        return p

    # ── Cover page ────────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("INTERNAL AUDIT REPORT")
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = NAVY

    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(data.get("report_title", "Audit Report"))
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = BLUE

    doc.add_paragraph()

    opinion = data.get("overall_opinion", "Needs Improvement")
    opinion_color = RISK_COLORS_DOCX.get(opinion, BLUE)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"Overall Opinion: {opinion}")
    run.font.size  = Pt(14)
    run.font.bold  = True
    run.font.color.rgb = opinion_color

    doc.add_paragraph()
    for label, value in [
        ("Audit Period:", data.get("audit_period", datetime.now().strftime("%B %Y"))),
        ("Date of Report:", datetime.now().strftime("%d %B %Y")),
        ("Classification:", "CONFIDENTIAL"),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r1 = p.add_run(f"{label}  ")
        r1.font.bold  = True
        r1.font.size  = Pt(11)
        r1.font.color.rgb = NAVY
        r2 = p.add_run(value)
        r2.font.size  = Pt(11)

    doc.add_page_break()

    # ── Executive Summary ─────────────────────────────────────────────────────
    add_heading("1. Executive Summary")
    add_body(data.get("executive_summary", ""))

    # Finding count table
    from collections import Counter
    findings = data.get("findings", [])
    risk_counts = Counter(f.get("risk_rating", "Medium") for f in findings)

    doc.add_paragraph()
    add_body("Summary of Findings by Risk Rating:", bold=True)
    tbl = doc.add_table(rows=2, cols=5)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for col_i, rating in enumerate(["Critical", "High", "Medium", "Low", "Informational"]):
        hcell = tbl.cell(0, col_i)
        hcell.text = rating
        set_cell_bg(hcell, RISK_HEX.get(rating, "CCCCCC"))
        hcell.paragraphs[0].runs[0].font.bold  = True
        hcell.paragraphs[0].runs[0].font.color.rgb = WHITE
        hcell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

        dcell = tbl.cell(1, col_i)
        dcell.text = str(risk_counts.get(rating, 0))
        dcell.paragraphs[0].runs[0].font.bold = True
        dcell.paragraphs[0].runs[0].font.size = Pt(16)
        dcell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()

    # ── Introduction ──────────────────────────────────────────────────────────
    add_heading("2. Scope & Methodology")
    if data.get("audit_scope"):
        add_body("Audit Scope", bold=True)
        add_body(data["audit_scope"])
    if data.get("methodology"):
        add_body("Methodology", bold=True)
        add_body(data["methodology"])

    doc.add_page_break()

    # ── Findings ──────────────────────────────────────────────────────────────
    add_heading("3. Findings & Recommendations")

    for idx, finding in enumerate(findings, start=1):
        risk = finding.get("risk_rating", "Medium")
        risk_hex   = RISK_HEX.get(risk, "CCCCCC")
        risk_color = RISK_COLORS_DOCX.get(risk, BLUE)

        # Finding header
        p = doc.add_paragraph()
        run = p.add_run(f"Finding {idx}: {finding.get('title', '')}")
        run.font.size  = Pt(12)
        run.font.bold  = True
        run.font.color.rgb = NAVY
        p.paragraph_format.space_before = Pt(12)

        # Finding table
        tbl = doc.add_table(rows=7, cols=2)
        tbl.style = "Table Grid"
        widths = [Inches(1.8), Inches(5.5)]

        rows_data = [
            ("Reference",               finding.get("finding_ref", f"F{idx:03d}")),
            ("Risk Rating",             risk),
            ("Description",             finding.get("description", "")),
            ("Impact",                  finding.get("impact", "")),
            ("Root Cause",              finding.get("root_cause", "")),
            ("Recommendation",          finding.get("recommendation", "")),
            ("Regulatory Implications", finding.get("regulatory_implications", "")),
        ]

        for row_i, (label, value) in enumerate(rows_data):
            label_cell = tbl.cell(row_i, 0)
            value_cell = tbl.cell(row_i, 1)
            label_cell.text = label
            set_cell_bg(label_cell, "0D2B55")
            label_cell.paragraphs[0].runs[0].font.bold  = True
            label_cell.paragraphs[0].runs[0].font.color.rgb = WHITE
            label_cell.paragraphs[0].runs[0].font.size  = Pt(10)

            value_cell.text = value
            value_cell.paragraphs[0].runs[0].font.size  = Pt(10)

            if label == "Risk Rating":
                set_cell_bg(value_cell, risk_hex)
                value_cell.paragraphs[0].runs[0].font.bold  = True
                value_cell.paragraphs[0].runs[0].font.color.rgb = WHITE

            if label == "Recommendation":
                value_cell.paragraphs[0].runs[0].font.bold = True

            for cell in [label_cell, value_cell]:
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.LEFT

        # Management response & due date
        if finding.get("management_response") or finding.get("due_date"):
            p = doc.add_paragraph()
            if finding.get("due_date"):
                r = p.add_run(f"Target Date: {finding['due_date']}  ")
                r.font.bold  = True
                r.font.size  = Pt(10)
                r.font.color.rgb = NAVY
            if finding.get("management_response"):
                r = p.add_run("Management Response: ")
                r.font.bold = True
                r.font.size = Pt(10)
                r2 = p.add_run(finding["management_response"])
                r2.font.size = Pt(10)

        doc.add_paragraph()

    doc.add_page_break()

    # ── Conclusion ────────────────────────────────────────────────────────────
    add_heading("4. Conclusion")
    add_body(data.get("conclusion", ""))

    if data.get("follow_up_actions"):
        add_body("Follow-up Actions:", bold=True)
        for action in data["follow_up_actions"]:
            add_body(f"  ▸  {action}")

    # ── Appendix ──────────────────────────────────────────────────────────────
    doc.add_page_break()
    add_heading("Appendix: Regulatory Framework")
    jurs = ["Switzerland (FINMA)", "Singapore (MAS)", "Hong Kong (SFC/HKMA)",
            "Bahamas (SCB)", "European Union", "United Kingdom (FCA/PRA)"]
    add_body("This audit covered the following jurisdictions:", bold=True)
    for j in jurs:
        add_body(f"  •  {j}")

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename  = f"{output_dir}/Audit_Report_{timestamp}.docx"
    doc.save(filename)
    return filename


def generate_regulatory_framework_docx(data: dict, output_dir: str = "outputs") -> str:
    """Generate a Word document for the regulatory framework (Agent 1 output)."""
    from docx import Document
    from docx.shared import Pt, RGBColor as DocxRGB, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin   = Cm(2.5)
        section.right_margin  = Cm(2.5)

    NAVY = DocxRGB(0x0D, 0x2B, 0x55)
    BLUE = DocxRGB(0x1A, 0x5F, 0x9E)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("REGULATORY FRAMEWORK")
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(data.get("audit_topic", ""))
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BLUE

    doc.add_paragraph(f"Date: {datetime.now().strftime('%d %B %Y')}")
    doc.add_page_break()

    for jur_data in data.get("jurisdictions", []):
        p = doc.add_paragraph()
        run = p.add_run(f"{jur_data.get('jurisdiction', '')}  —  {jur_data.get('regulator', '')}")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = NAVY
        p.paragraph_format.space_before = Pt(12)

        for reg in jur_data.get("key_regulations", []):
            p = doc.add_paragraph()
            run = p.add_run(reg.get("name", ""))
            run.font.bold = True
            run.font.size = Pt(11)
            run.font.color.rgb = BLUE

            p = doc.add_paragraph(reg.get("description", ""))
            p.runs[0].font.size = Pt(10)

            for req in reg.get("key_requirements", []):
                p = doc.add_paragraph(f"  •  {req}", style="List Bullet")
                p.runs[0].font.size = Pt(10)

        if jur_data.get("recent_changes"):
            p = doc.add_paragraph()
            run = p.add_run("Recent / Upcoming Changes: ")
            run.font.bold = True
            run.font.size = Pt(10)
            p.add_run(jur_data["recent_changes"]).font.size = Pt(10)

        doc.add_paragraph()

    if data.get("cross_jurisdictional_considerations"):
        p = doc.add_paragraph()
        run = p.add_run("Cross-Jurisdictional Considerations")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = NAVY
        doc.add_paragraph(data["cross_jurisdictional_considerations"])

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename  = f"{output_dir}/Regulatory_Framework_{timestamp}.docx"
    doc.save(filename)
    return filename


# ── Excel exports ──────────────────────────────────────────────────────────────

def generate_risk_analysis_excel(data: dict, output_dir: str = "outputs") -> str:
    """Tab 1 — Risk matrix + Regulations + Recommendations in 3 sheets."""
    import openpyxl
    from openpyxl.styles import PatternFill, Font, Alignment
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    HDR_FONT = Font(bold=True, color="FFFFFF", size=11)
    HDR_FILL = PatternFill("solid", fgColor="2D54D4")
    LEVEL_FILL = {
        "Critical": PatternFill("solid", fgColor="3D0A0A"),
        "High":     PatternFill("solid", fgColor="3D1A0A"),
        "Moderate": PatternFill("solid", fgColor="3D320A"),
    }

    def _hdr(ws, headers, col_widths):
        ws.append(headers)
        for cell in ws[1]:
            cell.font = HDR_FONT
            cell.fill = HDR_FILL
            cell.alignment = Alignment(wrap_text=True, vertical="center")
        ws.auto_filter.ref = ws.dimensions
        for i, w in enumerate(col_widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    def _row_align(ws):
        for cell in ws[ws.max_row]:
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    # Sheet 1: Risk Matrix
    ws1 = wb.active
    ws1.title = "Risk Matrix"
    _hdr(ws1, ["#", "Risk", "Level", "Description", "Impact", "Likelihood", "Expected Control"],
         [4, 28, 12, 42, 32, 12, 38])
    for i, r in enumerate(data.get("risks", []), 1):
        lv = r.get("level", "")
        ws1.append([i, r.get("name",""), lv, r.get("description",""), r.get("impact",""), r.get("likelihood",""), r.get("control","")])
        if lv in LEVEL_FILL:
            for cell in ws1[ws1.max_row]:
                cell.fill = LEVEL_FILL[lv]
        _row_align(ws1)

    # Sheet 2: Regulations
    ws2 = wb.create_sheet("Applicable Regulations")
    _hdr(ws2, ["Jurisdiction", "Regulation", "Reference", "Key Requirement"], [20, 35, 20, 60])
    for r in data.get("regs", []):
        ws2.append([r.get("jurisdiction",""), r.get("text",""), r.get("reference",""), r.get("requirement","")])
        _row_align(ws2)

    # Sheet 3: Public Recommendations
    ws3 = wb.create_sheet("Public Recommendations")
    _hdr(ws3, ["Source", "Year", "Theme", "Recommendation", "Priority"], [18, 8, 20, 60, 12])
    for r in data.get("pub_recs", []):
        ws3.append([r.get("source",""), r.get("year",""), r.get("theme",""), r.get("recommendation",""), r.get("priority","")])
        _row_align(ws3)

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{output_dir}/Risk_Analysis_{ts}.xlsx"
    wb.save(path)
    return path


def generate_audit_findings_excel(data: dict, output_dir: str = "outputs") -> str:
    """Tab 3 — Findings + Action Plan in 2 sheets."""
    import openpyxl
    from openpyxl.styles import PatternFill, Font, Alignment
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    HDR_FONT = Font(bold=True, color="FFFFFF", size=11)
    HDR_FILL = PatternFill("solid", fgColor="2D54D4")
    RATING_FILL = {
        "Critical": PatternFill("solid", fgColor="3D0A0A"),
        "High":     PatternFill("solid", fgColor="3D1A0A"),
        "Moderate": PatternFill("solid", fgColor="3D320A"),
        "Low":      PatternFill("solid", fgColor="0A3D1A"),
    }

    def _hdr(ws, headers, col_widths):
        ws.append(headers)
        for cell in ws[1]:
            cell.font = HDR_FONT; cell.fill = HDR_FILL
            cell.alignment = Alignment(wrap_text=True, vertical="center")
        ws.auto_filter.ref = ws.dimensions
        for i, w in enumerate(col_widths, 1):
            ws.column_dimensions[get_column_letter(i)].width = w

    ws1 = wb.active
    ws1.title = "Findings"
    _hdr(ws1, ["#", "Finding", "Rating", "Observation", "Risk", "Recommendation", "Owner", "Due Date"],
         [4, 30, 12, 45, 32, 45, 20, 14])
    for i, f in enumerate(data.get("findings", []), 1):
        rating = f.get("rating", f.get("severity", ""))
        ws1.append([i, f.get("title",""), rating, f.get("observation", f.get("description","")),
                    f.get("risk", f.get("impact","")), f.get("recommendation",""), f.get("owner",""), f.get("due_date","")])
        if rating in RATING_FILL:
            for cell in ws1[ws1.max_row]:
                cell.fill = RATING_FILL[rating]
        for cell in ws1[ws1.max_row]:
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    ws2 = wb.create_sheet("Action Plan")
    _hdr(ws2, ["#", "Finding", "Rating", "Action", "Owner", "Due Date", "Status"],
         [4, 30, 12, 55, 20, 14, 12])
    for i, f in enumerate(data.get("findings", []), 1):
        ws2.append([i, f.get("title",""), f.get("rating", f.get("severity","")),
                    f.get("recommendation",""), f.get("owner",""), f.get("due_date",""), "Open"])
        for cell in ws2[ws2.max_row]:
            cell.alignment = Alignment(wrap_text=True, vertical="top")

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{output_dir}/Audit_Findings_{ts}.xlsx"
    wb.save(path)
    return path


# ── Additional PowerPoint exports ──────────────────────────────────────────────

def generate_tab1_pptx(data: dict, output_dir: str = "outputs") -> str:
    """Tab 1 — Jurisdiction slides + risk heatmap."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BLUE  = RGBColor(0xA0, 0xB4, 0xF8)
    GRAY  = RGBColor(0x5A, 0x64, 0x88)
    BODY  = RGBColor(0xC8, 0xD0, 0xE8)
    RATING_RGB = {
        "Critical": RGBColor(0xEF, 0x44, 0x44),
        "High":     RGBColor(0xF9, 0x73, 0x16),
        "Moderate": RGBColor(0xEA, 0xB3, 0x08),
    }

    def _slide(title_txt, body_lines, title_color=None):
        s = prs.slides.add_slide(blank)
        t = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
        tf = t.text_frame
        tf.paragraphs[0].text = title_txt
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = title_color or BLUE
        b = s.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.9))
        btf = b.text_frame; btf.word_wrap = True
        for j, line in enumerate(body_lines):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            p.text = str(line); p.font.size = Pt(12.5)
            p.font.color.rgb = BODY
        return s

    topic = data.get("topic", "Audit Topic")
    risks = data.get("risks", []) or []
    regs  = data.get("regs", []) or []
    jurs  = data.get("jurs", []) or []

    # Title slide
    s0 = prs.slides.add_slide(blank)
    t0 = s0.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(1.3))
    t0.text_frame.paragraphs[0].text = f"Risk Analysis — {topic}"
    t0.text_frame.paragraphs[0].font.size = Pt(32)
    t0.text_frame.paragraphs[0].font.bold = True
    t0.text_frame.paragraphs[0].font.color.rgb = BLUE
    t1 = s0.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(0.6))
    t1.text_frame.paragraphs[0].text = "AuditIQ · Private Banking · Multi-Jurisdiction"
    t1.text_frame.paragraphs[0].font.size = Pt(14)
    t1.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Risk heatmap
    n_c = sum(1 for r in risks if r.get("level") == "Critical")
    n_h = sum(1 for r in risks if r.get("level") == "High")
    n_m = sum(1 for r in risks if r.get("level") == "Moderate")
    hm = [f"Total risks identified: {len(risks)}", "",
          f"  🔴 Critical : {n_c}", f"  🟠 High     : {n_h}", f"  🟡 Moderate : {n_m}", "",
          "Top Critical Risks:"] + [f"  • {r.get('name','')}" for r in risks if r.get("level") == "Critical"][:5]
    _slide(f"Risk Heatmap — {topic}", hm)

    # Per-jurisdiction slides
    for jur in jurs:
        jur_regs = [r for r in regs if r.get("jurisdiction","") == jur]
        lines = [f"Applicable regulations: {len(jur_regs)}", ""]
        for r in jur_regs[:8]:
            lines.append(f"  • {r.get('text','')} — {r.get('reference','')}")
        if len(jur_regs) > 8:
            lines.append(f"  … and {len(jur_regs)-8} more")
        _slide(f"Regulatory Landscape — {jur}", lines)

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{output_dir}/Risk_Analysis_{ts}.pptx"
    prs.save(path)
    return path


def generate_report_pptx(data: dict, output_dir: str = "outputs") -> str:
    """Tab 3 — Executive summary + one slide per finding."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    BLUE  = RGBColor(0xA0, 0xB4, 0xF8)
    GRAY  = RGBColor(0x5A, 0x64, 0x88)
    BODY  = RGBColor(0xC8, 0xD0, 0xE8)
    RATING_RGB = {
        "Critical": RGBColor(0xEF, 0x44, 0x44),
        "High":     RGBColor(0xF9, 0x73, 0x16),
        "Moderate": RGBColor(0xEA, 0xB3, 0x08),
        "Low":      RGBColor(0x22, 0xD3, 0xA5),
    }

    name     = data.get("name", "Audit Report")
    findings = data.get("findings", []) or []

    # Title slide
    s0 = prs.slides.add_slide(blank)
    t0 = s0.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.4))
    t0.text_frame.paragraphs[0].text = name
    t0.text_frame.paragraphs[0].font.size = Pt(28)
    t0.text_frame.paragraphs[0].font.bold = True
    t0.text_frame.paragraphs[0].font.color.rgb = BLUE
    t1 = s0.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(0.6))
    t1.text_frame.paragraphs[0].text = "AuditIQ · Internal Audit Report · Private Banking"
    t1.text_frame.paragraphs[0].font.size = Pt(14)
    t1.text_frame.paragraphs[0].font.color.rgb = GRAY

    # Executive summary
    n_hi = sum(1 for f in findings if f.get("rating","").lower() in ("critical","high"))
    exec_lines = [f"Total findings: {len(findings)}", f"Critical / High: {n_hi}", "",
                  "Key findings:"] + [f"  • [{f.get('rating','')}] {f.get('title','')}" for f in findings[:6]]
    s1 = prs.slides.add_slide(blank)
    th = s1.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.8))
    th.text_frame.paragraphs[0].text = "Executive Summary"
    th.text_frame.paragraphs[0].font.size = Pt(22); th.text_frame.paragraphs[0].font.bold = True
    th.text_frame.paragraphs[0].font.color.rgb = BLUE
    bx = s1.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.9))
    btf = bx.text_frame; btf.word_wrap = True
    for j, line in enumerate(exec_lines):
        p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
        p.text = line; p.font.size = Pt(13); p.font.color.rgb = BODY

    # Per-finding slides
    for f in findings:
        rating = f.get("rating", f.get("severity",""))
        color  = RATING_RGB.get(rating, RGBColor(0x8A, 0x92, 0xBB))
        s = prs.slides.add_slide(blank)
        th2 = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.75))
        th2.text_frame.paragraphs[0].text = f"Finding — {rating} — {f.get('title','')}"
        th2.text_frame.paragraphs[0].font.size = Pt(18)
        th2.text_frame.paragraphs[0].font.bold = True
        th2.text_frame.paragraphs[0].font.color.rgb = color
        content = [
            f"Observation: {f.get('observation', f.get('description',''))}",
            "",
            f"Risk: {f.get('risk', f.get('impact',''))}",
            "",
            f"Recommendation: {f.get('recommendation','')}",
            "",
            f"Owner: {f.get('owner','')}   Due: {f.get('due_date','')}",
        ]
        bx2 = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(6))
        btf2 = bx2.text_frame; btf2.word_wrap = True
        for j, line in enumerate(content):
            p = btf2.paragraphs[0] if j == 0 else btf2.add_paragraph()
            p.text = line; p.font.size = Pt(13); p.font.color.rgb = BODY

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"{output_dir}/Audit_Report_{ts}.pptx"
    prs.save(path)
    return path
